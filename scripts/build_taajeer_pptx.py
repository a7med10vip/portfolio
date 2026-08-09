#!/usr/bin/env python3
"""
Build Taajeer_Automotive.pptx from the web deck's own spec.

The web page and this file are generated from the same JSON (`/taajeer/spec`), so
the PowerPoint is the web deck — same coordinates, same Calibri, same colours —
and every element lands as a native, editable shape: rectangles, text boxes and
pictures. Nothing is flattened to an image.

The canvas is 1920x1080 px == 20in x 11.25in at 96 px/in, so:
    1 px = 9525 EMU   and   1 px = 0.75 pt
Both are exact — no rounding anywhere in the mapping.

Usage:
    npm run dev                       # the spec route must be served
    python3 scripts/build_taajeer_pptx.py [out.pptx]
"""

import copy
import io
import json
import sys
import urllib.request
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parent.parent
PUBLIC = ROOT / "public"
SPEC_URL = "http://localhost:3000/taajeer/spec"

EMU_PX = 9525          # 1 px
PT_PX = 0.75           # 1 px in points

ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCHOR = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def family(ff, fw):
    """Cambria has no Light cut, so the serif falls back to Regular."""
    if ff == "c":
        return "Cambria"
    return "Calibri Light" if fw == "l" else "Calibri"


def px(v):
    return Emu(int(round(v * EMU_PX)))


def pt(v):
    return Pt(v * PT_PX)


def rgb(hexstr):
    return RGBColor.from_string(hexstr.lstrip("#").upper()[:6])


def alpha_of(hexstr):
    """#RRGGBBAA -> alpha in PowerPoint's 0..100000, or None when fully opaque.

    The deck leans on translucent colour (ghost numerals, the cover scrim), so
    dropping the alpha channel here would silently turn a whisper into a shout.
    """
    h = hexstr.lstrip("#")
    if len(h) != 8:
        return None
    a = int(h[6:8], 16)
    return int(round(a / 255 * 100000))


def apply_alpha(color_format, hexstr):
    """Attach an <a:alpha> to a colour that python-pptx has already written."""
    a = alpha_of(hexstr)
    if a is None or a >= 100000:
        return
    clr = color_format._xFill.find(qn("a:srgbClr"))
    if clr is None:
        return
    old = clr.find(qn("a:alpha"))
    if old is not None:
        clr.remove(old)
    node = clr.makeelement(qn("a:alpha"), {"val": str(a)})
    clr.append(node)


def no_line(shape):
    shape.line.fill.background()


def sub(parent, tag, **attrs):
    el = parent.makeelement(qn(tag), {k: str(v) for k, v in attrs.items()})
    parent.append(el)
    return el


def soft_shadow(shape):
    """The card shadow, matched to the CSS: 0 6px 20px rgba(14,17,23,.06)."""
    spPr = shape._element.spPr
    for tag in ("a:effectLst",):
        old = spPr.find(qn(tag))
        if old is not None:
            spPr.remove(old)
    fx = sub(spPr, "a:effectLst")
    shdw = sub(fx, "a:outerShdw", blurRad=190500, dist=57150, dir=5400000, rotWithShape="0")
    clr = sub(shdw, "a:srgbClr", val="0E1117")
    sub(clr, "a:alpha", val="10000")


def round_picture(pic, rad, w, h):
    """Crop a picture to a rounded rectangle — PowerPoint's picture-to-shape."""
    spPr = pic._element.spPr
    old = spPr.find(qn("a:prstGeom"))
    if old is not None:
        spPr.remove(old)
    geom = spPr.makeelement(qn("a:prstGeom"), {"prst": "roundRect"})
    av = sub(geom, "a:avLst")
    # roundRect's adj is the radius as a fraction of the shorter side
    adj = min(rad / min(w, h), 0.5) if min(w, h) else 0
    sub(av, "a:gd", name="adj", fmla=f"val {int(adj * 100000)}")
    spPr.insert(0, geom)


# ── primitives ───────────────────────────────────────────────────────────────

def add_rect(slide, el):
    rad = el.get("rad") or 0
    if rad > 0:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(el["x"]), px(el["y"]), px(el["w"]), px(el["h"]))
        # PowerPoint's roundRect adjustment is the corner radius as a fraction of
        # the shorter side, capped at 0.5 (a pill).
        adj = min(rad / min(el["w"], el["h"]), 0.5)
        shp.adjustments[0] = adj
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(el["x"]), px(el["y"]), px(el["w"]), px(el["h"]))

    if el.get("fill"):
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(el["fill"])
        apply_alpha(shp.fill.fore_color, el["fill"])
    else:
        shp.fill.background()

    if el.get("line"):
        shp.line.color.rgb = rgb(el["line"])
        shp.line.width = pt(el.get("lw", 1))
        apply_alpha(shp.line.color, el["line"])
    else:
        no_line(shp)

    shp.shadow.inherit = False          # never take the theme's preset shadow
    if el.get("shadow"):
        soft_shadow(shp)
    shp.text_frame.text = ""
    return shp


def add_text(slide, el):
    box = slide.shapes.add_textbox(px(el["x"]), px(el["y"]), px(el["w"]), px(el["h"]))
    tf = box.text_frame
    tf.word_wrap = not el.get("nowrap", False)
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = ANCHOR[el.get("va", "t")]

    # autofit off — the box is already sized by the spec
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit"):
        node = bodyPr.find(qn(tag))
        if node is not None:
            bodyPr.remove(node)

    fs = el["fs"]
    fw = el.get("fw", "r")
    ff = el.get("ff", "s")
    it = el.get("it", False)
    base = el["c"]
    caps = el.get("caps", False)
    lh = el.get("lh")
    ls = el.get("ls")

    s = el["s"]
    runs_per_para = []
    if isinstance(s, str):
        for line in s.split("\n"):
            runs_per_para.append([{"t": line}])
    else:
        # runs may themselves contain hard breaks; split them into paragraphs
        para = []
        for rn in s:
            parts = rn["t"].split("\n")
            for k, part in enumerate(parts):
                if k > 0:
                    runs_per_para.append(para)
                    para = []
                nr = copy.deepcopy(rn)
                nr["t"] = part
                para.append(nr)
        runs_per_para.append(para)

    for pi, runs in enumerate(runs_per_para):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = ALIGN[el.get("al", "l")]
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        if lh:
            p.line_spacing = pt(lh)     # exact point spacing == CSS line-height in px
        for rn in runs:
            r = p.add_run()
            r.text = rn["t"].upper() if caps else rn["t"]
            rff = rn.get("ff", ff)          # a run may switch family mid-line
            f = r.font
            f.size = pt(fs)
            f.name = family(rff, fw)
            f.bold = fw == "b" or bool(rn.get("b"))
            f.italic = it or bool(rn.get("i"))
            rc = rn.get("c") or base
            f.color.rgb = rgb(rc)
            apply_alpha(f.color, rc)
            if ls:
                # a:rPr/@spc is tracking in 1/100 pt — python-pptx has no wrapper
                r._r.get_or_add_rPr().set("spc", str(int(round(ls * PT_PX * 100))))
    return box


def add_img(slide, el, missing):
    src = (PUBLIC / el["src"].lstrip("/")).resolve()
    if not src.exists():
        missing.add(el["src"])
        return None

    x, y, w, h = el["x"], el["y"], el["w"], el["h"]
    fit = el.get("fit", "contain")

    if fit == "contain":
        # letterbox by hand so the picture sits exactly where the browser puts it
        with Image.open(src) as im:
            iw, ih = im.size
        scale = min(w / iw, h / ih)
        dw, dh = iw * scale, ih * scale
        al = el.get("al", "c")
        dx = x if al == "l" else (x + w - dw if al == "r" else x + (w - dw) / 2)
        dy = y + (h - dh) / 2
        pic = slide.shapes.add_picture(str(src), px(dx), px(dy), px(dw), px(dh))
        if el.get("rad"):
            round_picture(pic, el["rad"], dw, dh)
        return pic

    # fit=cover: centre-crop to the box's aspect, the way object-fit:cover does
    pic = slide.shapes.add_picture(str(src), px(x), px(y), px(w), px(h))
    with Image.open(src) as im:
        iw, ih = im.size
    box_ar, img_ar = w / h, iw / ih
    if img_ar > box_ar:                      # image is wider — trim the sides
        frac = (1 - box_ar / img_ar) / 2
        pic.crop_left = pic.crop_right = frac
    elif img_ar < box_ar:                    # image is taller — trim top/bottom
        frac = (1 - img_ar / box_ar) / 2
        pic.crop_top = pic.crop_bottom = frac
    if el.get("rad"):
        round_picture(pic, el["rad"], w, h)
    return pic


# ── build ────────────────────────────────────────────────────────────────────

def build(spec, out):
    prs = Presentation()
    prs.slide_width = px(spec["w"])
    prs.slide_height = px(spec["h"])
    # python-pptx leaves the default 4:3 flag on sldSz even after resizing, which
    # makes PowerPoint offer to "fix" the aspect ratio on open.
    prs._element.find(qn("p:sldSz")).set("type", "custom")
    blank = prs.slide_layouts[6]

    missing = set()
    for s in spec["slides"]:
        slide = prs.slides.add_slide(blank)

        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = rgb(s.get("bg") or "#FFFFFF")

        for el in s["els"]:
            k = el["k"]
            if k == "r":
                add_rect(slide, el)
            elif k == "t":
                add_text(slide, el)
            elif k == "i":
                add_img(slide, el, missing)

    prs.save(out)
    return missing


def main():
    out = sys.argv[1] if len(sys.argv) > 1 else str(ROOT / "public" / "taajeer" / "Taajeer_Automotive.pptx")
    Path(out).parent.mkdir(parents=True, exist_ok=True)

    try:
        with urllib.request.urlopen(SPEC_URL, timeout=20) as r:
            spec = json.load(r)
    except Exception as e:
        sys.exit(f"could not read {SPEC_URL} — is `npm run dev` running?\n  {e}")

    missing = build(spec, out)
    print(f"✓ {len(spec['slides'])} slides → {out}")
    if missing:
        print("\n⚠ images referenced by the spec but not on disk (slides built without them):")
        for m in sorted(missing):
            print("   ", m)


if __name__ == "__main__":
    main()
